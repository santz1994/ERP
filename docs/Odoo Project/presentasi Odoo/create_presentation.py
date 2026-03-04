from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def add_content_slide(prs, title, bullet_points=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if bullet_points:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(14)
    
    return slide

def add_blank_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    return slide

# Title Slide
title_slide = add_title_slide(
    prs,
    "PT Quty Karunia",
    "Evaluasi & Rencana Implementasi Odoo ERP 2026\nSoft Toys & Stuffing Manufacturing | IKEA Supplier"
)

# Slide 1 - Overview Workflow Manual
slide1 = add_blank_slide(prs)
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 1 — Workflow Saat Ini (Manual Process)"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(74, 144, 217)

content_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

sections = [
    "⚠️ KONDISI SAAT INI:",
    "Seluruh proses operasional PT Quty Karunia masih berbasis Excel manual dan komunikasi verbal.",
    "Purchasing menjadi trigger utama material planning — bukan PPIC.",
    "",
    "ALUR PURCHASING → WAREHOUSE (Manual):",
    "• Terima Order → Buka Excel BOM (478 SKU)",
    "• Kalkulasi Manual Kebutuhan Material",
    "• Cek Stock di Excel Gudang (Delay 1-2 jam)",
    "• Buat PO di Excel (Copy-paste, rawan typo)",
    "• Email PO ke Supplier (Manual)",
    "• Warehouse Terima Material",
    "• Admin Update Excel Stock (End-of-day, rawan salah)",
]

for i, section in enumerate(sections):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = section
    p.font.size = Pt(14)
    if "⚠️" in section or "KONDISI" in section or "ALUR" in section:
        p.font.bold = True
        p.font.size = Pt(16)

# Slide 1b - Pain Points
slide1b = add_blank_slide(prs)
title_box = slide1b.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Pain Points Utama - Proses Manual"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(208, 2, 27)

content_box = slide1b.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

pain_points = [
    "1. Material shortage tiba-tiba → Production stop, overtime, emergency purchasing",
    "2. Tidak tahu WIP real-time → Planning impossible, bottleneck detection delayed",
    "3. QC defect manual (Excel terpisah) → Tidak terintegrasi, IKEA audit risk",
    "4. 🚨 CRITICAL: Mix Label risk (Finishing-Packing) → IKEA REJECT!",
    "5. Dual trigger system manual → Admin overwhelmed, packaging error frequent",
    "6. Excel-based planning → Error prone, single point of failure",
    "7. Target produksi flat → Shortage frequent (tidak hitung defect rate)",
    "8. Manual material tracking → Admin time wasted 2-3 jam/hari",
    "9. Stock opname 1 hari penuh → Production stop untuk counting",
    "10. Rework tidak ter-record → Cost tidak tahu, IKEA audit concern",
    "11. Multi-unit conversion manual + Pallet → UOM error: ROLL↔METER, KG↔GRAM",
]

for i, point in enumerate(pain_points):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = point
    p.font.size = Pt(13)
    if "CRITICAL" in point:
        p.font.bold = True
        p.font.color.rgb = RGBColor(208, 2, 27)

# Slide 2 - Odoo 2026 Workflow
slide2 = add_blank_slide(prs)
title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 2 — Workflow Rencana Odoo 2026"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(126, 211, 33)

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

odoo_content = [
    "✅ VISI 2026:",
    "Seluruh proses terintegrasi dalam Odoo Enterprise.",
    "Dari input order hingga dispatch — otomatis, real-time, dan zero mix label.",
    "",
    "ALUR PURCHASING → WAREHOUSE (Odoo Otomatis):",
    "• Input Order di Odoo → ODOO AUTOMATION",
    "• Auto Explode BOM (478 SKU instant!)",
    "• Auto Convert Units ROLL→PCS, KG→GRAM",
    "• IKEA Pallet Logic: Karton Genap per Pallet",
    "• Purchase Requisition List (Auto-generated)",
    "• 1-Click Create PO + Email Auto-send ke Supplier!",
    "• Warehouse Scan Mobile + Validasi Barcode",
    "• Stock Update REAL-TIME!",
    "• Dashboard: Material Ready → Info ke Produksi: VALID",
]

for i, section in enumerate(odoo_content):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = section
    p.font.size = Pt(14)
    if "✅" in section or "VISI" in section or "ALUR" in section:
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(126, 211, 33)

# Slide 2b - Manfaat Odoo 2026
slide2b = add_blank_slide(prs)
title_box = slide2b.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Manfaat Utama Odoo 2026"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Three columns for benefits
col_width = 2.8
left_box = slide2b.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(col_width), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

purchasing = [
    "🛒 PURCHASING",
    "• Auto-calculate material dari BOM",
    "• Multi-unit conversion otomatis",
    "• IKEA pallet rules built-in",
    "• Stock visibility real-time",
    "• PO tracking otomatis",
]

for i, item in enumerate(purchasing):
    p = left_frame.add_paragraph() if i > 0 else left_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(12)
    if "🛒" in item:
        p.font.bold = True
        p.font.size = Pt(14)

mid_box = slide2b.shapes.add_textbox(Inches(3.6), Inches(1.2), Inches(col_width), Inches(5))
mid_frame = mid_box.text_frame
mid_frame.word_wrap = True

production = [
    "🏭 PRODUCTION",
    "• MO & WO auto-generate",
    "• Material backflush otomatis",
    "• QC terintegrasi per dept",
    "• Rework tracked & traceable",
    "• 38 sewing lines tracking otomatis",
]

for i, item in enumerate(production):
    p = mid_frame.add_paragraph() if i > 0 else mid_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(12)
    if "🏭" in item:
        p.font.bold = True
        p.font.size = Pt(14)

right_box = slide2b.shapes.add_textbox(Inches(6.7), Inches(1.2), Inches(col_width), Inches(5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

management = [
    "📊 MANAGEMENT",
    "• Real-time WIP dashboard",
    "• Bottleneck alert otomatis",
    "• Defect pattern analysis",
    "• IKEA audit compliance ready",
    "• Completion forecast akurat",
]

for i, item in enumerate(management):
    p = right_frame.add_paragraph() if i > 0 else right_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(12)
    if "📊" in item:
        p.font.bold = True
        p.font.size = Pt(14)

# Slide 3 - Comparison Manual vs Odoo
slide3 = add_blank_slide(prs)
title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 3 — Perbandingan: Manual vs Odoo 2026"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

content_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

comparison = [
    "KALKULASI MATERIAL:",
    "🔴 Manual: Excel & Rumus. Rawan salah konversi unit",
    "🟢 Odoo: Otomatis. System hitung kebutuhan + konversi + pallet genap",
    "",
    "STOCK CHECK:",
    "🔴 Manual: Excel + hubungi Warehouse (delay 1-2 jam)",
    "🟢 Odoo: REAL-TIME, No Delay",
    "",
    "CREATE PO:",
    "🔴 Manual: Excel (copy-paste rawan typo)",
    "🟢 Odoo: Auto-generate + Email auto-send!",
    "",
    "LABEL / DATESTAMP:",
    "🔴 Manual: Manual Look-up. RISIKO TINGGI Mix Label!",
    "🟢 Odoo: Auto-Inherit dari PO Label. Tidak bisa diedit manual.",
]

for i, item in enumerate(comparison):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    if ":" in item and not "🔴" in item and not "🟢" in item:
        p.font.bold = True
        p.font.size = Pt(14)

# Slide 4 - Odoo 2023 Gagal
slide4 = add_blank_slide(prs)
title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 4 — Odoo 2023: Apa yang Gagal?"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(208, 2, 27)

content_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

failures = [
    "💀 Project Odoo 2023 (Odoo Community Edition via Impact) — GAGAL TOTAL",
    "Proyek di-abandon setelah 6 bulan, tim kembali ke Excel.",
    "",
    "8 MASALAH KRITIS:",
    "❌ 1. BOM Management: Input manual 478 SKU → 3-4 bulan data entry",
    "❌ 2. Material Availability: Status MERAH palsu → Production block",
    "❌ 3. Reporting: Tidak ada export → Data trapped di system",
    "❌ 4. Vendor Force-fit: 'Kalian yang harus adjust' → User frustasi",
    "❌ 5. UOM Conversion: Mismatch ROLL vs METER → BOM consumption salah",
    "❌ 6. UI/UX: Menu 5-6 klik, Bahasa Inggris → User kembali ke Excel",
    "❌ 7. MO Manual: 20-30 MO/hari manual create → Human error",
    "❌ 8. WO Manual: 5 dept = 5x manual work → WO tidak sync",
]

for i, item in enumerate(failures):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    if "💀" in item or "8 MASALAH" in item:
        p.font.bold = True
        p.font.size = Pt(15)

# Slide 5 - Odoo 2026 Berbeda
slide5 = add_blank_slide(prs)
title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 5 — Odoo 2026: Apa yang Berbeda?"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(126, 211, 33)

content_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

solutions = [
    "🚀 Rencana 2026: Odoo Enterprise via Odoo Indonesia (Principal)",
    "Pendekatan partnership, bukan transactional.",
    "",
    "8 SOLUSI UNTUK 8 MASALAH:",
    "✅ 1. BOM: Import massal 478 SKU + Dual BOM (Purchasing + Production)",
    "✅ 2. Material: Dual MO System + Real-time stock sync",
    "✅ 3. Reporting: Export Excel/CSV + API integration + Dashboard",
    "✅ 4. Partnership: Discovery workshop + Gap analysis transparent",
    "✅ 5. UOM: Multi-unit configured + Pallet-level tracking (IKEA rules)",
    "✅ 6. UI/UX: Simplified menu (1-2 klik) + Bahasa Indonesia + Mobile",
    "✅ 7. MO: PO → Auto-generate 2 MO parallel (IKEA + Pabrik)",
    "✅ 8. WO: 1 MO → 5 WO/SPK auto-generated + Real-time visibility",
]

for i, item in enumerate(solutions):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    if "🚀" in item or "8 SOLUSI" in item:
        p.font.bold = True
        p.font.size = Pt(15)

# Slide 5b - Dual MO System
slide5b = add_blank_slide(prs)
title_box = slide5b.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Dual MO System — Rekomendasi ⭐"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

content_box = slide5b.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

dual_mo = [
    "📋 PO Purchasing → 2 MO Parallel:",
    "",
    "1️⃣ MO PURCHASING (Fixed Planning)",
    "   • Target: 1000 pcs (Fixed)",
    "   • Untuk: IKEA Reporting",
    "   • Status: Clean, Match Order",
    "",
    "2️⃣ MO PABRIK (Dynamic Execution)",
    "   • Target: Flexible",
    "   • Untuk: Factory Tracking",
    "   • Input: Daily Production + Consumption + QC",
    "",
    "🔄 RECONCILIATION OTOMATIS:",
    "   Variance: +50 pcs (5% efficiency gain) → GREEN ✅",
]

for i, item in enumerate(dual_mo):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(14)
    if "📋" in item or "🔄" in item or "1️⃣" in item or "2️⃣" in item:
        p.font.bold = True

# Slide 6 - Comparison 2023 vs 2026
slide6 = add_blank_slide(prs)
title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 6 — Perbandingan Lengkap: 2023 vs 2026"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

content_box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

comparison_2 = [
    "ODOO COMMUNITY vs ODOO ENTERPRISE:",
    "🔴 Community: Free, Fitur Terbatas, Community Support, Manual Upgrade",
    "🟢 Enterprise: Berbayar, Fitur Lengkap, Official Support, Auto Upgrade",
    "",
    "IMPLEMENTASI 2023 vs 2026:",
    "Setup BOM:",
    "  🔴 2023: Input manual satu per satu",
    "  🟢 2026: Import massal + Dual BOM",
    "",
    "Sistem MO:",
    "  🔴 2023: Single MO (variance membingungkan)",
    "  🟢 2026: Dual MO (IKEA tetap + Pabrik dinamis) ⭐ RECOMMENDED",
    "",
    "UI/UX:",
    "  🔴 2023: Kompleks, hanya Bahasa Inggris",
    "  🟢 2026: Disederhanakan, Bahasa Indonesia",
]

for i, item in enumerate(comparison_2):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    if ":" in item and not "🔴" in item and not "🟢" in item and not " " in item[:2]:
        p.font.bold = True
        p.font.size = Pt(14)

# Slide 7 - Odoo Indonesia & Metodologi
slide7 = add_blank_slide(prs)
title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Slide 7 — Metodologi Implementasi Odoo Indonesia"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

content_box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

methodology = [
    "🏢 ODOO INDONESIA (Principal) — Representasi Resmi Odoo S.A.",
    "",
    "TAHAPAN IMPLEMENTASI:",
    "1️⃣ INISIASI (Scoping): Kirim requirements & pain points",
    "2️⃣ GAP ANALYSIS: Consultation & Estimasi",
    "3️⃣ PELAKSANAAN: Deep dive operasional PT Quty",
    "4️⃣ DELIVERABLES: Gap Analysis Report → Validasi → Revisi jika perlu",
    "5️⃣ FINALISASI: Proposal Proyek Final",
    "",
    "JAMINAN KESESUAIAN:",
    "🔍 Filter 1: Product Demo sebelum development",
    "✅ Filter 2: User Acceptance Testing (UAT)",
    "🏥 Hypercare: 2-3 bulan setelah Go-Live",
    "",
    "🤝 PARTNERSHIP APPROACH:",
    "'Kami demand partnership, bukan transactional.'",
]

for i, item in enumerate(methodology):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    if any(emoji in item for emoji in ["🏢", "TAHAPAN", "JAMINAN", "🤝"]):
        p.font.bold = True
        p.font.size = Pt(14)

# Slide 8 - Kesimpulan
slide8 = add_blank_slide(prs)
title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Kesimpulan"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(74, 144, 217)

content_box = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

conclusion = [
    "🎯 KESIMPULAN:",
    "",
    "Implementasi Odoo 2026 dirancang berdasarkan lessons learned",
    "dari kegagalan 2023.",
    "",
    "Dengan pendekatan:",
    "✅ Dual BOM + Dual MO System",
    "✅ 7 unique requirements ter-address",
    "✅ Partnership approach dengan Odoo Indonesia (Principal)",
    "✅ Enterprise Edition dengan fitur lengkap",
    "✅ Metodologi terstruktur dengan jaminan kesesuaian",
    "",
    "PT Quty Karunia siap untuk transformasi digital yang sukses!",
]

for i, item in enumerate(conclusion):
    p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    if "🎯" in item:
        p.font.bold = True
        p.font.size = Pt(24)
    if "sukses!" in item:
        p.font.bold = True
        p.font.color.rgb = RGBColor(126, 211, 33)

# Save presentation
prs.save('PT_Quty_Karunia_Odoo_2026.pptx')
print("✅ Presentasi berhasil dibuat: PT_Quty_Karunia_Odoo_2026.pptx")
