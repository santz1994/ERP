#!/usr/bin/env python3
"""
Convert PRESENTASI_MANAGEMENT_ERP_QUTY_KARUNIA.md to DOCX and PPTX formats
"""

import re
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN

def clean_markdown_text(text):
    """Remove markdown formatting from text"""
    # Remove markdown links [text](url)
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    # Remove bold **text**
    text = re.sub(r'\*\*([^\*]+)\*\*', r'\1', text)
    # Remove italic *text*
    text = re.sub(r'\*([^\*]+)\*', r'\1', text)
    # Remove code blocks ```
    text = re.sub(r'```[^\n]*\n(.*?)```', r'\1', text, flags=re.DOTALL)
    # Remove inline code `text`
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Remove emojis (basic)
    text = re.sub(r'[\U0001F300-\U0001F9FF]', '', text)
    return text

def convert_to_docx(md_file_path, output_path):
    """Convert markdown to DOCX"""
    print(f"Converting {md_file_path} to DOCX...")
    
    # Read markdown file
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Create document
    doc = Document()
    
    # Split content by lines
    lines = content.split('\n')
    
    in_code_block = False
    code_block_content = []
    in_table = False
    
    for line in lines:
        # Handle code blocks
        if line.startswith('```'):
            if in_code_block:
                # End of code block
                if code_block_content:
                    code_text = '\n'.join(code_block_content)
                    p = doc.add_paragraph(code_text)
                    p.style = 'No Spacing'
                    for run in p.runs:
                        run.font.name = 'Courier New'
                        run.font.size = Pt(9)
                code_block_content = []
                in_code_block = False
            else:
                # Start of code block
                in_code_block = True
            continue
        
        if in_code_block:
            code_block_content.append(line)
            continue
        
        # Skip horizontal rules
        if line.strip() == '---':
            doc.add_paragraph()
            continue
        
        # Skip table separators
        if re.match(r'^\|[\s\-:]+\|', line):
            continue
        
        # Handle tables
        if line.startswith('|') and '|' in line[1:]:
            # Simple table handling - just add as paragraph
            clean_line = clean_markdown_text(line.strip('|').strip())
            if clean_line:
                p = doc.add_paragraph(clean_line)
                p.style = 'No Spacing'
            continue
        
        # Handle headings
        if line.startswith('#'):
            level = len(line) - len(line.lstrip('#'))
            heading_text = clean_markdown_text(line.lstrip('#').strip())
            if heading_text:
                h = doc.add_heading(heading_text, level=min(level, 9))
            continue
        
        # Handle list items
        if re.match(r'^\s*[-*]\s+', line):
            text = clean_markdown_text(re.sub(r'^\s*[-*]\s+', '', line))
            if text:
                doc.add_paragraph(text, style='List Bullet')
            continue
        
        # Handle numbered lists
        if re.match(r'^\s*\d+\.\s+', line):
            text = clean_markdown_text(re.sub(r'^\s*\d+\.\s+', '', line))
            if text:
                doc.add_paragraph(text, style='List Number')
            continue
        
        # Regular paragraphs
        if line.strip():
            text = clean_markdown_text(line.strip())
            if text and not text.startswith('<'):  # Skip HTML comments
                doc.add_paragraph(text)
    
    # Save document
    doc.save(output_path)
    print(f"✅ DOCX saved to: {output_path}")

def convert_to_pptx(md_file_path, output_path):
    """Convert markdown to PPTX"""
    print(f"Converting {md_file_path} to PPTX...")
    
    # Read markdown file
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)
    
    # Split content into sections by ## headings
    sections = re.split(r'\n## ', content)
    
    # First section is title
    if sections:
        title_section = sections[0]
        title_lines = title_section.split('\n')
        title = clean_markdown_text(title_lines[0].lstrip('#').strip())
        subtitle = clean_markdown_text(title_lines[1].strip()) if len(title_lines) > 1 else ""
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle
    
    # Process remaining sections
    for section in sections[1:]:
        lines = section.split('\n')
        if not lines:
            continue
        
        # Section heading becomes slide title
        slide_title = clean_markdown_text(lines[0].strip())
        
        # Use bullet layout
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = slide_title
        
        # Content body
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            
            bullet_count = 0
            for line in lines[1:]:
                line = line.strip()
                
                # Skip empty lines, horizontal rules, code blocks
                if not line or line == '---' or line.startswith('```'):
                    continue
                
                # Skip HTML comments and anchors
                if line.startswith('<') or line.startswith('|'):
                    continue
                
                # Handle list items
                if re.match(r'^[-*]\s+', line):
                    text = clean_markdown_text(re.sub(r'^[-*]\s+', '', line))
                    if text and bullet_count < 15:  # Limit bullets per slide
                        p = text_frame.add_paragraph()
                        p.text = text
                        p.level = 0
                        bullet_count += 1
                    continue
                
                # Handle numbered lists
                if re.match(r'^\d+\.\s+', line):
                    text = clean_markdown_text(re.sub(r'^\d+\.\s+', '', line))
                    if text and bullet_count < 15:
                        p = text_frame.add_paragraph()
                        p.text = text
                        p.level = 0
                        bullet_count += 1
                    continue
                
                # Handle subheadings
                if line.startswith('###'):
                    heading = clean_markdown_text(line.lstrip('#').strip())
                    if heading and bullet_count < 15:
                        p = text_frame.add_paragraph()
                        p.text = heading
                        p.level = 0
                        p.font.bold = True
                        bullet_count += 1
                    continue
                
                # Regular text (limit to avoid overflow)
                if bullet_count < 15:
                    text = clean_markdown_text(line)
                    if text and len(text) < 200:
                        p = text_frame.add_paragraph()
                        p.text = text[:150]  # Truncate long text
                        p.level = 1
                        bullet_count += 1
    
    # Save presentation
    prs.save(output_path)
    print(f"✅ PPTX saved to: {output_path}")

def main():
    md_file = "docs/00-Overview/PRESENTASI_MANAGEMENT_ERP_QUTY_KARUNIA.md"
    
    # Output paths
    docx_output = "docs/00-Overview/PRESENTASI_MANAGEMENT_ERP_QUTY_KARUNIA.docx"
    pptx_output = "docs/00-Overview/PRESENTASI_MANAGEMENT_ERP_QUTY_KARUNIA.pptx"
    
    # Convert to DOCX
    convert_to_docx(md_file, docx_output)
    
    # Convert to PPTX
    convert_to_pptx(md_file, pptx_output)
    
    print("\n✅ Conversion complete!")
    print(f"📄 DOCX: {docx_output}")
    print(f"📊 PPTX: {pptx_output}")

if __name__ == "__main__":
    main()
